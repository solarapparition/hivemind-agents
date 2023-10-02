from pptx import Presentation

# Create a PowerPoint presentation
presentation = Presentation()

# Function to add slides based on the YAML content
def add_slide(title, content):
    slide_layout = presentation.slide_layouts[1]  # Title and Content layout
    slide = presentation.slides.add_slide(slide_layout)
    title_box = slide.shapes.title
    content_box = slide.placeholders[1]
    title_box.text = title
    for item in content:
        content_box.text += f"\u2022 {item}\n"

# Iterate through the YAML content and create slides
for section in yaml_content['sections']:
    add_slide(section['Title'], section['Content'])

# Save the PowerPoint file
pptx_path = '/mnt/data/output.pptx'
presentation.save(pptx_path)

pptx_path
